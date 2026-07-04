from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from engine.content import ContentItem
from engine.image_resolver import resolve_image
from engine.layout_selector import choose_layout


def _hex(color: str) -> RGBColor:
    color = color.lstrip("#")
    return RGBColor.from_string(color)


def _blank_slide(prs, theme, bg_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _hex(bg_color or theme.get("background", "FFFFFF"))
    return slide


def _caption(slide, theme, source):
    if not source:
        return
    cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9), Inches(0.4))
    cap.text_frame.text = f"Source: {source}"
    cap.text_frame.paragraphs[0].font.size = Pt(10)
    cap.text_frame.paragraphs[0].font.color.rgb = _hex(theme.get("muted_color", "888888"))


# ---------------------------------------------------------------------------
# Layout implementations. Each takes (slide, theme, item) and draws content.
# Add new ones here and register them in LAYOUT_REGISTRY below.
# ---------------------------------------------------------------------------

def _layout_two_column_image(slide, theme, item: ContentItem):
    image_path = resolve_image(item.image_path)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = item.title or "Untitled"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _hex(theme["primary_color"])
    tf.paragraphs[0].font.name = theme.get("font_header", "Cambria")

    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.2), Inches(4.8))
    if image_path:
        slide.shapes.add_picture(image_path, Inches(6.0), Inches(1.4), width=Inches(3.5))

    body_tf = text_box.text_frame
    body_tf.word_wrap = True
    body_tf.text = item.text
    body_tf.paragraphs[0].font.size = Pt(15)
    body_tf.paragraphs[0].font.color.rgb = _hex(theme.get("text_color", "212121"))
    body_tf.paragraphs[0].font.name = theme.get("font_body", "Calibri")
    body_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    _caption(slide, theme, item.source)


def _layout_full_bleed_image(slide, theme, item: ContentItem):
    image_path = resolve_image(item.image_path)
    if image_path:
        # Fill left 55% of slide with the image
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), height=Inches(7.5), width=Inches(5.5))
        text_left = Inches(6.0)
        text_width = Inches(3.5)
    else:
        text_left = Inches(0.5)
        text_width = Inches(9)

    title_box = slide.shapes.add_textbox(text_left, Inches(0.5), text_width, Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = item.title or "Untitled"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _hex(theme["primary_color"])
    tf.paragraphs[0].font.name = theme.get("font_header", "Cambria")

    body_box = slide.shapes.add_textbox(text_left, Inches(1.6), text_width, Inches(5.0))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.text = item.text
    body_tf.paragraphs[0].font.size = Pt(14)
    body_tf.paragraphs[0].font.color.rgb = _hex(theme.get("text_color", "212121"))
    body_tf.paragraphs[0].font.name = theme.get("font_body", "Calibri")

    _caption(slide, theme, item.source)


def _layout_text_block(slide, theme, item: ContentItem):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = item.title or "Untitled"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _hex(theme["primary_color"])
    tf.paragraphs[0].font.name = theme.get("font_header", "Cambria")

    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.2))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.text = item.text
    body_tf.paragraphs[0].font.size = Pt(16)
    body_tf.paragraphs[0].font.color.rgb = _hex(theme.get("text_color", "212121"))
    body_tf.paragraphs[0].font.name = theme.get("font_body", "Calibri")
    body_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    _caption(slide, theme, item.source)


def _layout_stat_callout(slide, theme, item: ContentItem):
    # Big centered number/short text, title as a small label above it
    label_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(0.6))
    ltf = label_box.text_frame
    ltf.text = (item.title or "").upper()
    ltf.paragraphs[0].font.size = Pt(16)
    ltf.paragraphs[0].font.color.rgb = _hex(theme.get("muted_color", "888888"))
    ltf.paragraphs[0].font.name = theme.get("font_body", "Calibri")
    ltf.paragraphs[0].alignment = PP_ALIGN.CENTER

    stat_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9.0), Inches(1.8))
    stf = stat_box.text_frame
    stf.word_wrap = True
    stf.text = item.text
    stf.paragraphs[0].font.size = Pt(60)
    stf.paragraphs[0].font.bold = True
    stf.paragraphs[0].font.color.rgb = _hex(theme["primary_color"])
    stf.paragraphs[0].font.name = theme.get("font_header", "Cambria")
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER

    _caption(slide, theme, item.source)


def _layout_quote(slide, theme, item: ContentItem):
    quote_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(7.6), Inches(2.2))
    qtf = quote_box.text_frame
    qtf.word_wrap = True
    qtf.text = f"\u201c{item.text}\u201d"
    qtf.paragraphs[0].font.size = Pt(26)
    qtf.paragraphs[0].font.italic = True
    qtf.paragraphs[0].font.color.rgb = _hex(theme["primary_color"])
    qtf.paragraphs[0].font.name = theme.get("font_header", "Cambria")
    qtf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if item.title:
        attr_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.9), Inches(7.6), Inches(0.5))
        atf = attr_box.text_frame
        atf.text = f"— {item.title}"
        atf.paragraphs[0].font.size = Pt(14)
        atf.paragraphs[0].font.color.rgb = _hex(theme.get("muted_color", "888888"))
        atf.paragraphs[0].font.name = theme.get("font_body", "Calibri")
        atf.paragraphs[0].alignment = PP_ALIGN.CENTER

    _caption(slide, theme, item.source)


def _layout_image_grid(slide, theme, item: ContentItem):
    all_images = ([item.image_path] if item.image_path else []) + list(item.image_paths)
    resolved = [resolve_image(p) for p in all_images]
    resolved = [p for p in resolved if p][:4]  # cap at 4

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = item.title or "Untitled"
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _hex(theme["primary_color"])
    tf.paragraphs[0].font.name = theme.get("font_header", "Cambria")

    # 2x2 grid positions within a 9" x 4.6" area starting at (0.5, 1.3)
    cell_w, cell_h = Inches(4.4), Inches(2.2)
    gap = Inches(0.2)
    positions = [
        (Inches(0.5), Inches(1.3)),
        (Inches(0.5) + cell_w + gap, Inches(1.3)),
        (Inches(0.5), Inches(1.3) + cell_h + gap),
        (Inches(0.5) + cell_w + gap, Inches(1.3) + cell_h + gap),
    ]
    for path, (left, top) in zip(resolved, positions):
        slide.shapes.add_picture(path, left, top, width=cell_w, height=cell_h)

    if item.text:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.6))
        ctf = cap_box.text_frame
        ctf.word_wrap = True
        ctf.text = item.text
        ctf.paragraphs[0].font.size = Pt(12)
        ctf.paragraphs[0].font.color.rgb = _hex(theme.get("text_color", "212121"))
        ctf.paragraphs[0].font.name = theme.get("font_body", "Calibri")

    _caption(slide, theme, item.source)


def _layout_bulleted_list(slide, theme, item: ContentItem):
    import re as _re

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = item.title or "Untitled"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _hex(theme["primary_color"])
    tf.paragraphs[0].font.name = theme.get("font_header", "Cambria")

    # Split on newlines, semicolons, or leading bullet markers
    raw_items = _re.split(r"\n|;", item.text)
    bullets = [b.strip(" -•\t") for b in raw_items if b.strip(" -•\t")]

    body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.0))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = _hex(theme.get("text_color", "212121"))
        p.font.name = theme.get("font_body", "Calibri")
        p.space_after = Pt(10)

    _caption(slide, theme, item.source)


LAYOUT_REGISTRY = {
    "two_column_image": _layout_two_column_image,
    "full_bleed_image": _layout_full_bleed_image,
    "text_block": _layout_text_block,
    "stat_callout": _layout_stat_callout,
    "quote": _layout_quote,
    "image_grid": _layout_image_grid,
    "bulleted_list": _layout_bulleted_list,
}


def _choose_layout(item: ContentItem, template: dict, index: int) -> str:
    return choose_layout(item, template, index, set(LAYOUT_REGISTRY.keys()))


def _add_title_slide(prs, theme, deck_title, subtitle=""):
    slide = _blank_slide(prs, theme, bg_color=theme["primary_color"])

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.5))
    tf = tb.text_frame
    tf.text = deck_title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = _hex(theme.get("on_primary", "FFFFFF"))
    p.font.name = theme.get("font_header", "Cambria")

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(8.4), Inches(0.8))
        tf2 = tb2.text_frame
        tf2.text = subtitle
        tf2.paragraphs[0].font.size = Pt(18)
        tf2.paragraphs[0].font.color.rgb = _hex(theme.get("secondary_color", "CADCFC"))
        tf2.paragraphs[0].font.name = theme.get("font_body", "Calibri")
    return slide


def render_pptx(
    items: list[ContentItem],
    template: dict,
    deck_title: str,
    output_path: str,
    subtitle: str = "",
):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    theme = template["theme"]

    _add_title_slide(prs, theme, deck_title, subtitle)

    for index, item in enumerate(items):
        layout_name = _choose_layout(item, template, index)
        layout_fn = LAYOUT_REGISTRY[layout_name]
        slide = _blank_slide(prs, theme)
        layout_fn(slide, theme, item)

    prs.save(output_path)
    return output_path
